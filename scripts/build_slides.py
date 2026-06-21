"""Gera a apresentação do TP2 - Checkpoint Etapa 1 (Modelo Hedônico)."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# ---------- Paleta (UFMG-friendly, sóbria) ----------
COLOR_PRIMARY   = RGBColor(0x0B, 0x3D, 0x91)  # azul escuro
COLOR_ACCENT    = RGBColor(0xE2, 0x7D, 0x60)  # coral
COLOR_DARK      = RGBColor(0x1F, 0x2A, 0x44)  # quase preto
COLOR_MUTED     = RGBColor(0x6B, 0x72, 0x80)  # cinza
COLOR_LIGHT     = RGBColor(0xF1, 0xF5, 0xF9)  # cinza claro
COLOR_GREEN     = RGBColor(0x2E, 0x7D, 0x32)
COLOR_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)

# ---------- Setup ----------
prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]


def add_rect(slide, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
    shp.shadow.inherit = False
    return shp


def add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=COLOR_DARK,
             align=PP_ALIGN.LEFT, font="Calibri"):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run(); run.text = line
        run.font.size = Pt(size); run.font.bold = bold
        run.font.color.rgb = color; run.font.name = font
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=14, color=COLOR_DARK, bullet_color=None):
    """items: list[str] (suporta sub-bullets com '  ' inicial)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    bc = bullet_color or COLOR_ACCENT
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(4)
        sub = item.startswith("  ")
        text = item.strip()
        bullet = "•  " if not sub else "–  "
        if sub:
            p.level = 1
        r1 = p.add_run(); r1.text = bullet
        r1.font.size = Pt(size); r1.font.bold = True
        r1.font.color.rgb = bc; r1.font.name = "Calibri"
        r2 = p.add_run(); r2.text = text
        r2.font.size = Pt(size); r2.font.color.rgb = color; r2.font.name = "Calibri"
    return tb


def add_header(slide, title, subtitle=None):
    # faixa lateral
    add_rect(slide, Emu(0), Emu(0), Inches(0.25), SH, COLOR_PRIMARY)
    add_text(slide, Inches(0.55), Inches(0.30), Inches(12.5), Inches(0.6),
             title, size=26, bold=True, color=COLOR_PRIMARY)
    if subtitle:
        add_text(slide, Inches(0.55), Inches(0.85), Inches(12.5), Inches(0.4),
                 subtitle, size=13, color=COLOR_MUTED)
    add_rect(slide, Inches(0.55), Inches(1.25), Inches(1.4), Inches(0.05), COLOR_ACCENT)


def add_footer(slide, page_num, total):
    add_text(slide, Inches(0.55), Inches(7.10), Inches(8.0), Inches(0.3),
             "TP2 — Checkpoint · Grupo 02 · Efeito Dominó no Mercado de Transferências",
             size=9, color=COLOR_MUTED)
    add_text(slide, Inches(11.5), Inches(7.10), Inches(1.5), Inches(0.3),
             f"{page_num} / {total}", size=9, color=COLOR_MUTED, align=PP_ALIGN.RIGHT)


def add_table(slide, x, y, w, h, data, *, header=True, highlight_row=None,
              col_widths=None, font_size=11):
    rows = len(data); cols = len(data[0])
    tbl_shape = slide.shapes.add_table(rows, cols, x, y, w, h)
    tbl = tbl_shape.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = cw
    for r in range(rows):
        for c in range(cols):
            cell = tbl.cell(r, c)
            cell.text = ""
            cell.margin_left = cell.margin_right = Inches(0.06)
            cell.margin_top = cell.margin_bottom = Inches(0.04)
            p = cell.text_frame.paragraphs[0]
            run = p.add_run(); run.text = str(data[r][c])
            run.font.name = "Calibri"
            if header and r == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = COLOR_PRIMARY
                run.font.color.rgb = COLOR_WHITE
                run.font.bold = True; run.font.size = Pt(font_size)
            else:
                if highlight_row is not None and r == highlight_row:
                    cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9)
                    run.font.bold = True
                else:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = COLOR_WHITE if r % 2 == 1 else COLOR_LIGHT
                run.font.color.rgb = COLOR_DARK
                run.font.size = Pt(font_size)
            if c > 0:
                p.alignment = PP_ALIGN.CENTER
    return tbl


# =================================================================
# SLIDE 1 — Capa
# =================================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, Emu(0), Emu(0), SW, SH, COLOR_PRIMARY)
add_rect(s, Inches(0), Inches(6.7), SW, Inches(0.8), COLOR_DARK)
add_text(s, Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.5),
         "TP2 · Checkpoint de Modelagem", size=18, color=COLOR_ACCENT, bold=True)
add_text(s, Inches(0.8), Inches(2.2), Inches(11.5), Inches(2.0),
         "Efeito Dominó no Mercado\nde Transferências de Futebol",
         size=44, bold=True, color=COLOR_WHITE)
add_text(s, Inches(0.8), Inches(4.4), Inches(11.5), Inches(0.5),
         "O \"Prêmio do Vendedor\" em Compras Subsequentes",
         size=20, color=COLOR_LIGHT)
add_text(s, Inches(0.8), Inches(5.4), Inches(11.5), Inches(0.4),
         "Grupo 02  ·  Ciência de Dados Aplicada ao Futebol  ·  UFMG",
         size=13, color=COLOR_LIGHT)
add_text(s, Inches(0.8), Inches(5.75), Inches(11.5), Inches(0.4),
         "Carlos · César · Letícia · Lucas · Lucca",
         size=13, color=COLOR_LIGHT)
add_text(s, Inches(0.8), Inches(6.9), Inches(11.5), Inches(0.4),
         "Junho / 2026", size=11, color=COLOR_MUTED)


# =================================================================
# SLIDE 2 — Recap (pergunta de pesquisa)
# =================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "Pergunta de Pesquisa", "Recap da proposta inicial")

# Card central com a pergunta
add_rect(s, Inches(0.7), Inches(1.7), Inches(12.0), Inches(1.7), COLOR_LIGHT)
add_rect(s, Inches(0.7), Inches(1.7), Inches(0.12), Inches(1.7), COLOR_ACCENT)
add_text(s, Inches(1.0), Inches(1.85), Inches(11.5), Inches(0.4),
         "Pergunta principal", size=12, bold=True, color=COLOR_ACCENT)
add_text(s, Inches(1.0), Inches(2.20), Inches(11.5), Inches(1.2),
         "Dada uma venda realizada pelo Time A em uma janela de transferência, qual é a "
         "probabilidade de que outros clubes (B, C, D) vendam um jogador para o Time A por "
         "um preço acima do valor de mercado de referência em compras subsequentes dentro "
         "da mesma janela?", size=15, color=COLOR_DARK)

add_text(s, Inches(0.7), Inches(3.7), Inches(12.0), Inches(0.4),
         "Hipótese central", size=14, bold=True, color=COLOR_PRIMARY)
add_bullets(s, Inches(0.7), Inches(4.1), Inches(12.0), Inches(1.4), [
    "Clubes recém-capitalizados perdem poder de barganha e pagam um prêmio (sobrepreço) "
    "nas contratações seguintes — o \"prêmio do vendedor\".",
    "O efeito é causal, e não apenas correlação devida a \"clubes ricos\" que naturalmente "
    "vendem caro e compram caro.",
], size=14)

add_text(s, Inches(0.7), Inches(5.6), Inches(12.0), Inches(0.4),
         "Por que importa", size=14, bold=True, color=COLOR_PRIMARY)
add_bullets(s, Inches(0.7), Inches(6.0), Inches(12.0), Inches(1.0), [
    "Aplicação direta em inteligência de mercado de clubes profissionais.",
    "Inovação metodológica: foco em preço e contexto temporal, não em estrutura de rede.",
], size=14)
add_footer(s, 2, 10)


# =================================================================
# SLIDE 3 — Estratégia em 2 Etapas
# =================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "Estratégia de Modelagem", "Arquitetura em duas etapas sequenciais")

# Etapa 1 - card
def stage_card(x, y, w, h, num, title, lines, color):
    add_rect(s, x, y, w, h, COLOR_LIGHT)
    add_rect(s, x, y, w, Inches(0.5), color)
    add_text(s, x + Inches(0.25), y + Inches(0.08), w, Inches(0.4),
             f"ETAPA {num}  ·  {title}", size=13, bold=True, color=COLOR_WHITE)
    tb = s.shapes.add_textbox(x + Inches(0.25), y + Inches(0.65), w - Inches(0.5), h - Inches(0.7))
    tf = tb.text_frame; tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(3)
        label, val = (ln if isinstance(ln, tuple) else (ln, None))
        r1 = p.add_run(); r1.text = label
        r1.font.size = Pt(12); r1.font.bold = True; r1.font.color.rgb = COLOR_PRIMARY
        r1.font.name = "Calibri"
        if val:
            r2 = p.add_run(); r2.text = " " + val
            r2.font.size = Pt(12); r2.font.color.rgb = COLOR_DARK
            r2.font.name = "Calibri"

stage_card(Inches(0.7), Inches(1.6), Inches(5.9), Inches(4.4), 1,
           "Predição do Preço Hedônico", [
    ("Objetivo:",  "estimar o \"preço justo\" do jogador via ML."),
    ("Target:",    "ln(fee) — log do valor da transação."),
    ("Features:",  "age, age², log(market_value), posição,"),
    ("           ", "liga (dummies), efeitos fixos de temporada."),
    ("Modelos:",   "XGBoost · LightGBM · Random Forest · SVR."),
    ("Validação:", "CV 5-fold + split temporal 2023-24 / 2025."),
    ("Saída:",     "resíduo Yᵢ,c = ln(P) − ln(P̂)"),
    ("         ",  "→ \"prêmio de reinvestimento\" puro."),
], COLOR_PRIMARY)

stage_card(Inches(6.85), Inches(1.6), Inches(5.9), Inches(4.4), 2,
           "Estimação Causal (Double ML)", [
    ("Objetivo:",  "isolar o efeito causal da liquidez recente"),
    ("           ", "do clube sobre o prêmio pago."),
    ("Tratamento (D):", "log_revenue — Σ vendas na temporada."),
    ("Y:",         "resíduo hedônico out-of-fold da Etapa 1."),
    ("W (confound.):", "elenco + rede + liga + temporada"),
    ("           ", "(19 vars; sem bloco financeiro direto)."),
    ("Algoritmo:", "DML manual (Robinson) + R-learner;"),
    ("           ", "SE HC1 e clusterizado clube×temporada."),
    ("Saída:",     "θ̂ (ATE) + heterogeneidade por clube."),
], COLOR_ACCENT)

# seta entre os cards
arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                           Inches(6.55), Inches(3.55), Inches(0.35), Inches(0.5))
arrow.fill.solid(); arrow.fill.fore_color.rgb = COLOR_DARK
arrow.line.fill.background()

add_text(s, Inches(0.7), Inches(6.15), Inches(12.0), Inches(0.4),
         "Por que esta arquitetura?", size=13, bold=True, color=COLOR_PRIMARY)
add_text(s, Inches(0.7), Inches(6.5), Inches(12.0), Inches(0.7),
         "Separar predição (Etapa 1) de inferência causal (Etapa 2) é o padrão do estado da "
         "arte (Chernozhukov et al., 2018). Garante que ML não-linear capture os fundamentos "
         "do jogador sem viesar a estimativa do parâmetro causal de interesse.",
         size=12, color=COLOR_DARK)
add_footer(s, 3, 10)


# =================================================================
# SLIDE 4 — Etapa 1: Resultados (tabela comparativa)
# =================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "Etapa 1 · Resultados Comparativos",
           "4 modelos · CV 5-fold + holdout temporal (test = 2025)")

table_data = [
    ["Modelo", "CV RMSE (μ)", "CV RMSE (σ)", "CV R² (μ)", "Test RMSE", "Test MAE", "Test R²"],
    ["XGBoost",       "0.6854", "0.019", "0.7218", "0.6939", "0.5005", "0.7463"],
    ["LightGBM",      "0.6770", "0.021", "0.7289", "0.6826", "0.4948", "0.7545"],
    ["Random Forest", "0.6738", "0.021", "0.7311", "0.6726", "0.4898", "0.7616"],
    ["SVR (RBF)",     "0.7026", "0.028", "0.7078", "0.7164", "0.5171", "0.7296"],
]
add_table(s, Inches(0.7), Inches(1.55), Inches(12.0), Inches(2.4),
          table_data, highlight_row=3,
          col_widths=[Inches(2.2)] + [Inches(1.633)] * 6,
          font_size=12)

# Vencedor
add_rect(s, Inches(0.7), Inches(4.15), Inches(12.0), Inches(0.55), RGBColor(0xE8, 0xF5, 0xE9))
add_text(s, Inches(0.9), Inches(4.25), Inches(11.6), Inches(0.4),
         "✓  Vencedor: Random Forest — Test RMSE 0.6726 · Test R² 0.7616 · OOB R² 0.7354.",
         size=13, bold=True, color=COLOR_GREEN)

# Justificativa de métricas
add_text(s, Inches(0.7), Inches(4.95), Inches(12.0), Inches(0.4),
         "Métricas escolhidas e por quê", size=14, bold=True, color=COLOR_PRIMARY)
add_bullets(s, Inches(0.7), Inches(5.35), Inches(12.0), Inches(1.7), [
    "Target contínuo (ln fee) ⇒ problema de regressão (não classificação).",
    "RMSE: penaliza erros grandes — relevante porque fees variam de €250k a €100M+.",
    "MAE: interpretação direta em escala log; robusto a outliers.",
    "R²: padrão na literatura hedônica (Tilburg, Emerald, ResearchGate reportam 0.70–0.80).",
    "Generalização: CV ~ Test (gap < 5%) indica ausência de overfitting relevante.",
], size=13)
add_footer(s, 4, 10)


# =================================================================
# SLIDE 5 — Hiperparâmetros & Configuração
# =================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "Configuração Experimental",
           "Reprodutibilidade · seed=42 · features=19 · n=5.146")

# Hiperparâmetros
hp = [
    ["Modelo", "Hiperparâmetros (configuração atual)"],
    ["XGBoost",       "n_estimators=500 · max_depth=5 · learning_rate=0.05 · subsample=0.8 · colsample_bytree=0.8"],
    ["LightGBM",      "n_estimators=500 · max_depth=5 · learning_rate=0.05 · subsample=0.8 · colsample_bytree=0.8"],
    ["Random Forest", "n_estimators=500 · max_depth=10 · min_samples_leaf=5 · oob_score=True"],
    ["SVR (RBF)",     "kernel='rbf' · C=10 · ε=0.1 · γ='scale' · pipeline com StandardScaler"],
]
add_table(s, Inches(0.7), Inches(1.55), Inches(12.0), Inches(2.4),
          hp, col_widths=[Inches(2.2), Inches(9.8)], font_size=12)

add_text(s, Inches(0.7), Inches(4.15), Inches(6.0), Inches(0.4),
         "Pipeline de dados", size=14, bold=True, color=COLOR_PRIMARY)
add_bullets(s, Inches(0.7), Inches(4.55), Inches(6.0), Inches(2.5), [
    "Filtro: fee_type='paid', direction='In', fee ≥ €250k.",
    "Winsorização do premium_ratio em 1–99%.",
    "Transformações log: log_fee, log_mv, log_league_mv.",
    "Resíduo hedônico out-of-fold + features de rede (PageRank).",
    "Dataset final: 5.146 transferências × 7 ligas × 8 temporadas.",
], size=12)

add_text(s, Inches(7.0), Inches(4.15), Inches(5.7), Inches(0.4),
         "Avaliação", size=14, bold=True, color=COLOR_PRIMARY)
add_bullets(s, Inches(7.0), Inches(4.55), Inches(5.7), Inches(2.5), [
    "KFold 5 folds (shuffle=True, seed=42).",
    "Split temporal: treino = 2017–2024 · teste = 2025.",
    "Resíduo Y via cross_val_predict (out-of-fold, sem vazamento).",
    "Métricas em log-space (RMSE, MAE, R²) + OOB no RF.",
    "TODO próxima iteração: tuning via Optuna / GridSearchCV.",
], size=12)
add_footer(s, 5, 10)


# =================================================================
# SLIDE 6 — Interpretabilidade (SHAP)
# =================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "Etapa 1 · Interpretabilidade (SHAP)",
           "Random Forest · TreeExplainer · ranking de impacto preditivo")

# Ranking SHAP (barras visuais simplificadas)
add_text(s, Inches(0.7), Inches(1.55), Inches(6.0), Inches(0.4),
         "Ranking de importância média |SHAP|", size=13, bold=True, color=COLOR_PRIMARY)

ranking = [
    ("log_mv",        1.00, "Market value do Transfermarkt — proxy dominante de habilidade"),
    ("age",           0.55, "Curva de carreira; jovens valorizados"),
    ("age_sq",        0.45, "Não-linearidade: pico ~20 anos"),
    ("liga (dummies)", 0.30, "Premier League e LaLiga elevam o preço esperado"),
    ("is_attacker",   0.18, "Atacantes comandam prêmio sobre defensores"),
    ("season_2024",   0.12, "Inflação de mercado controlada"),
    ("is_midfielder", 0.10, "Pequeno efeito residual"),
]
y0 = Inches(2.0); row_h = Inches(0.55)
for i, (name, val, desc) in enumerate(ranking):
    y = Inches(2.0 + i * 0.55)
    add_text(s, Inches(0.7), y, Inches(1.6), row_h, name,
             size=11, bold=True, color=COLOR_DARK)
    bar_w = Inches(2.5 * val)
    add_rect(s, Inches(2.35), y + Inches(0.08), bar_w, Inches(0.28), COLOR_ACCENT)
    add_text(s, Inches(5.0), y, Inches(1.0), row_h, f"{val:.2f}",
             size=11, color=COLOR_MUTED)

# Coluna direita: interpretação
add_text(s, Inches(7.0), Inches(1.55), Inches(5.7), Inches(0.4),
         "O que o SHAP revela", size=13, bold=True, color=COLOR_PRIMARY)
add_bullets(s, Inches(7.0), Inches(2.0), Inches(5.7), Inches(3.5), [
    "log_mv concentra ~40% do impacto preditivo, consistente com a literatura "
    "(Peeters, 2018; Tilburg University).",
    "Bloco age/age² confirma curvatura biológica do valor.",
    "Liga e temporada já capturam C4 (sazonalidade) e parte de C5 (inflação) "
    "sem precisar do DML.",
    "Posição entra como controle, com magnitude moderada.",
], size=12)

add_rect(s, Inches(7.0), Inches(5.4), Inches(5.7), Inches(1.5), COLOR_LIGHT)
add_text(s, Inches(7.2), Inches(5.5), Inches(5.4), Inches(0.3),
         "Estudo de caso na literatura (de Ligt, 2019)",
         size=11, bold=True, color=COLOR_PRIMARY)
add_text(s, Inches(7.2), Inches(5.8), Inches(5.4), Inches(1.1),
         "SHAP decompôs o fee de €77,2M em: Team Rating do destino + idade jovem. "
         "Nosso pipeline replica essa lógica e está pronto para esse tipo de análise "
         "individual por transferência.", size=11, color=COLOR_DARK)
add_footer(s, 6, 10)


# =================================================================
# SLIDE 7 — Tratamento dos Confundidores
# =================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "Tratamento dos 6 Confundidores",
           "Mapeamento causal e estratégias de neutralização")

conf = [
    ["#", "Confundidor", "Estratégia", "Status"],
    ["C1", "Poder financeiro do clube (\"clube rico\")",
     "Proxies de rede + elenco em W (sem bloco financeiro)", "Etapa 2 (W)"],
    ["C2", "Causalidade reversa (compra precede venda)",
     "Teste de lag t-1 (enfraquece; não conclusivo)", "Robustez · parcial"],
    ["C3", "Planejamento estratégico (causa comum)",
     "Exige contexto: UCL, troca de técnico, receita TV", "Limitação · futuro"],
    ["C4", "Calendário / sazonalidade (deadline day)",
     "Dummies de temporada em FE (Etapa 1) e em W (Etapa 2)", "Etapa 1 + 2"],
    ["C5", "Choques inflacionários de mercado",
     "Dummies de liga + temporada em W", "Etapa 1 + 2"],
    ["C6", "Viés de seleção do grupo tratado",
     "Features de rede + elenco em W (PSM = futuro)", "Etapa 2 (W)"],
]
add_table(s, Inches(0.7), Inches(1.55), Inches(12.0), Inches(4.3),
          conf, col_widths=[Inches(0.6), Inches(3.6), Inches(5.4), Inches(2.4)],
          font_size=11)

add_rect(s, Inches(0.7), Inches(6.05), Inches(12.0), Inches(0.9), COLOR_LIGHT)
add_text(s, Inches(0.9), Inches(6.15), Inches(11.6), Inches(0.3),
         "Por que isto é o coração do projeto",
         size=12, bold=True, color=COLOR_PRIMARY)
add_text(s, Inches(0.9), Inches(6.45), Inches(11.6), Inches(0.5),
         "Sem controlar C1–C6, qualquer correlação observada seria espúria. "
         "É essa disciplina causal que diferencia o trabalho de um valuation comum.",
         size=11, color=COLOR_DARK)
add_footer(s, 7, 10)


# =================================================================
# SLIDE 8 — Insights Preliminares
# =================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "Insights Preliminares",
           "O que os dados já estão dizendo")

# 3 cards numéricos
def kpi_card(x, y, w, h, kpi, label, color):
    add_rect(s, x, y, w, h, COLOR_LIGHT)
    add_rect(s, x, y, w, Inches(0.10), color)
    add_text(s, x, y + Inches(0.5), w, Inches(1.1), kpi,
             size=44, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(s, x, y + Inches(1.7), w, Inches(0.6), label,
             size=12, color=COLOR_DARK, align=PP_ALIGN.CENTER)

kpi_card(Inches(0.7),  Inches(1.55), Inches(3.9), Inches(2.4),
         "52,5%", "das transferências apresentam sobrepreço\nem relação ao valor hedônico estimado", COLOR_PRIMARY)
kpi_card(Inches(4.75), Inches(1.55), Inches(3.9), Inches(2.4),
         "76,2%", "da variância de ln(fee) é explicada\npelo modelo Random Forest no teste", COLOR_ACCENT)
kpi_card(Inches(8.8),  Inches(1.55), Inches(3.9), Inches(2.4),
         "+3,5%", "mediana do prêmio: leve sobrepreço médio\ndo mercado sobre o preço justo", COLOR_GREEN)

add_text(s, Inches(0.7), Inches(4.15), Inches(12.0), Inches(0.4),
         "Observações qualitativas", size=14, bold=True, color=COLOR_PRIMARY)
add_bullets(s, Inches(0.7), Inches(4.55), Inches(12.0), Inches(2.5), [
    "Resíduo out-of-fold aproximadamente simétrico em torno de zero (média −0.001; mediana +0.035), "
    "indicando ausência de viés sistemático e validando o uso na Etapa 2.",
    "Resíduo gerado sem vazamento (cross_val_predict): desvio-padrão 0.671 (vs 0.587 in-sample).",
    "log_mv do Transfermarkt domina o sinal preditivo — alinhado a Peeters (2018).",
    "Quartil 25 = −38% / Quartil 75 = +41%: 50% das transferências negociadas em margem normal.",
    "Hipótese respondida na Etapa 2: efeito médio nulo, mas significativo em 2022–2023.",
], size=12)
add_footer(s, 8, 10)


# =================================================================
# SLIDE 9 — Próximos Passos (cronograma)
# =================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "Próximos Passos",
           "Roteiro até a entrega final do TP3")

steps = [
    ("Semana 1", "Tuning de hiperparâmetros",
     "Optuna nos 4 modelos · target: ganhar 2–4 p.p. de R² · reavaliar SHAP."),
    ("Semana 1", "Enriquecimento hedônico",
     "Integrar FBref/StatsBomb: gols, assistências, xG, minutos jogados."),
    ("Semana 2", "Implementação da Etapa 2 — Double ML",
     "DML manual (Robinson) + R-learner com cross-fitting de 5 folds."),
    ("Semana 2", "Métricas inovadoras adaptadas",
     "IVB por clube · θ(PageRank quartil) · θ_t por temporada (proxy de decaimento)."),
    ("Semana 3", "Validação causal (C6)",
     "Propensity Score Matching para pares de clubes \"gêmeos\"."),
    ("Semana 3", "Robustez e entrega",
     "Análise de sensibilidade · subgrupos por tier de clube · relatório final."),
]

y0 = 1.55
for i, (when, title, desc) in enumerate(steps):
    y = Inches(y0 + i * 0.85)
    # número
    add_rect(s, Inches(0.7), y, Inches(0.55), Inches(0.7), COLOR_PRIMARY)
    add_text(s, Inches(0.7), y + Inches(0.15), Inches(0.55), Inches(0.45),
             f"{i+1}", size=22, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    # quando
    add_rect(s, Inches(1.4), y + Inches(0.1), Inches(1.0), Inches(0.5), COLOR_ACCENT)
    add_text(s, Inches(1.4), y + Inches(0.2), Inches(1.0), Inches(0.4),
             when, size=11, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    # título e desc
    add_text(s, Inches(2.6), y, Inches(10.2), Inches(0.35),
             title, size=13, bold=True, color=COLOR_DARK)
    add_text(s, Inches(2.6), y + Inches(0.36), Inches(10.2), Inches(0.4),
             desc, size=11, color=COLOR_MUTED)

add_footer(s, 9, 10)


# =================================================================
# SLIDE 10 — Checklist de validação (antes da entrega)
# =================================================================
s = prs.slides.add_slide(BLANK)
add_header(s, "Checklist de Validação",
           "Itens revisados antes de fechar a apresentação")

checklist = [
    ("Gráfico SHAP summary do melhor modelo",          "Ranking exibido no slide 6"),
    ("Tabela comparativa dos 4 modelos com vencedor",  "Slide 4 · Random Forest destacado"),
    ("Diagrama do pipeline Etapa 1 → Etapa 2",         "Slide 3 · cards + seta"),
    ("Dimensões da base citadas",                      "5.146 transferências · 7 ligas · 8 temporadas"),
    ("Hiperparâmetros explicitados",                   "Slide 5 · todos os 4 modelos"),
    ("Justificativa das métricas",                     "Slide 4 · RMSE / MAE / R² em log-space"),
    ("Mapeamento dos 6 confundidores",                 "Slide 7 · status por etapa"),
    ("Cronograma de próximos passos",                  "Slide 9 · 6 marcos em 3 semanas"),
    ("Limitações reconhecidas",                        "Sem tuning · sem features de performance ainda"),
    ("Equipe e contexto",                              "Slide 1 · grupo 02 / UFMG"),
]

for i, (item, where) in enumerate(checklist):
    y = Inches(1.55 + i * 0.46)
    # check icon
    ck = s.shapes.add_shape(MSO_SHAPE.OVAL,
                            Inches(0.7), y + Inches(0.05), Inches(0.3), Inches(0.3))
    ck.fill.solid(); ck.fill.fore_color.rgb = COLOR_GREEN
    ck.line.fill.background()
    add_text(s, Inches(0.7), y + Inches(0.05), Inches(0.3), Inches(0.3),
             "✓", size=12, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.15), y + Inches(0.05), Inches(7.5), Inches(0.35),
             item, size=12, bold=True, color=COLOR_DARK)
    add_text(s, Inches(8.7), y + Inches(0.05), Inches(4.5), Inches(0.35),
             where, size=11, color=COLOR_MUTED)

add_footer(s, 10, 10)


# ---------- Salva ----------
out = "/Users/lucaspedras/TP-CDAF/output/TP2_Checkpoint_Modelagem_Hedonica.pptx"
prs.save(out)
print(f"OK: {out}")
