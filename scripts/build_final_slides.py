"""Gera a apresentação FINAL do TP4 (DemoDay) — Grupo 02.

14 slides cobrindo a história completa: pergunta -> armadilha causal -> metodologia
(hedônico + Double ML) -> 3 achados -> robustez -> conclusões.

Fonte da verdade dos números: docs/script_apresentacao.md (alinhado aos notebooks).
Saída: output/TP4_Apresentacao_Final_Grupo02.pptx
"""
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# ---------- Paleta ----------
COLOR_PRIMARY = RGBColor(0x0B, 0x3D, 0x91)   # azul escuro
COLOR_ACCENT  = RGBColor(0xE2, 0x7D, 0x60)   # coral
COLOR_DARK    = RGBColor(0x1F, 0x2A, 0x44)
COLOR_MUTED   = RGBColor(0x6B, 0x72, 0x80)
COLOR_LIGHT   = RGBColor(0xF1, 0xF5, 0xF9)
COLOR_GREEN   = RGBColor(0x2E, 0x7D, 0x32)
COLOR_RED     = RGBColor(0xC0, 0x39, 0x2B)
COLOR_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
HEX_PRIMARY, HEX_ACCENT, HEX_GREEN, HEX_MUTED = "#0B3D91", "#E27D60", "#2E7D32", "#6B7280"

ASSETS = "/tmp/tp4_assets"
os.makedirs(ASSETS, exist_ok=True)

# ============================================================
# Gráficos (matplotlib) — números verificados nos notebooks
# ============================================================
def fig_progressao():
    """Slide 8: OLS ingênuo -> OLS+W -> DML, com IC (HC1)."""
    labels = ["OLS\nsem controles", "OLS\n+ W", "Double ML\n(não-linear)"]
    theta = [0.0142, 0.0051, 0.0044]
    lo    = [0.0092, -0.0001, -0.0016]
    hi    = [0.0192, 0.0103, 0.0104]
    colors = [HEX_ACCENT, HEX_MUTED, HEX_PRIMARY]
    fig, ax = plt.subplots(figsize=(7.2, 4.0), dpi=200)
    x = range(3)
    for i in x:
        ax.errorbar(i, theta[i], yerr=[[theta[i]-lo[i]], [hi[i]-theta[i]]],
                    fmt="o", ms=13, color=colors[i], ecolor="#9aa0a6",
                    elinewidth=2, capsize=6, capthick=2)
        sig = "*" if lo[i] > 0 else "n.s."
        ax.annotate(f"{theta[i]:.4f}\n{sig}", (i, theta[i]), xytext=(16, 0),
                    textcoords="offset points", va="center", fontsize=11,
                    fontweight="bold", color=colors[i])
    ax.axhline(0, color=HEX_RED if False else "#C0392B", ls="--", lw=1.2, alpha=.8)
    ax.set_xticks(list(x)); ax.set_xticklabels(labels, fontsize=10)
    ax.set_ylabel("θ (efeito da liquidez sobre o sobrepreço)", fontsize=10)
    ax.set_title("A correlação encolhe ao controlar o 'clube rico'", fontsize=12, fontweight="bold")
    ax.set_xlim(-0.4, 2.7); ax.grid(axis="y", alpha=.3)
    for s in ["top", "right"]: ax.spines[s].set_visible(False)
    fig.tight_layout(); p = f"{ASSETS}/progressao.png"; fig.savefig(p); plt.close(fig)
    return p

def fig_sazonal():
    """Slide 9: θ por temporada com IC clusterizado; 2022/2023 destacados."""
    seasons = [2017, 2018, 2019, 2021, 2022, 2023, 2024, 2025]
    theta = [-0.006, -0.013, 0.024, 0.012, 0.064, 0.047, 0.003, 0.015]
    lo    = [-0.026, -0.036, -0.011, -0.001, 0.028, 0.017, -0.053, -0.010]
    hi    = [0.014, 0.011, 0.058, 0.025, 0.101, 0.077, 0.059, 0.040]
    fig, ax = plt.subplots(figsize=(8.0, 4.2), dpi=200)
    for i, s in enumerate(seasons):
        sig = lo[i] > 0
        col = HEX_GREEN if sig else HEX_MUTED
        ax.errorbar(s, theta[i], yerr=[[theta[i]-lo[i]], [hi[i]-theta[i]]],
                    fmt="o", ms=10 if sig else 7, color=col, ecolor="#b0b4ba",
                    elinewidth=1.8, capsize=4, capthick=1.8, zorder=3 if sig else 2)
    ax.axhline(0, color="#C0392B", ls="--", lw=1.2, alpha=.8, label="θ = 0")
    ax.axvspan(2021.5, 2023.5, color="#E27D60", alpha=.10, zorder=0)
    ax.annotate("boom pós-COVID\n2022: +6,4%*  2023: +4,7%*\n(sobrevivem Bonferroni/FDR)",
                (2022.5, 0.101), xytext=(2018.6, 0.082), fontsize=9.5, color=HEX_GREEN,
                fontweight="bold", ha="left",
                arrowprops=dict(arrowstyle="->", color=HEX_GREEN, lw=1.3))
    ax.set_xticks(seasons); ax.set_ylabel("θ (sobrepreço causal)", fontsize=10)
    ax.set_title("Onde o efeito vive: significativo só em 2022–2023", fontsize=12, fontweight="bold")
    ax.grid(axis="y", alpha=.3); ax.legend(loc="lower right", fontsize=9)
    for sp in ["top", "right"]: ax.spines[sp].set_visible(False)
    fig.tight_layout(); p = f"{ASSETS}/sazonal.png"; fig.savefig(p); plt.close(fig)
    return p

IMG_PROG = fig_progressao()
IMG_SAZ  = fig_sazonal()

# ============================================================
# Helpers de layout
# ============================================================
prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def add_rect(s, x, y, w, h, fill, line=None):
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line if line else fill
    if not line: shp.line.fill.background()
    shp.shadow.inherit = False
    return shp

def add_text(s, x, y, w, h, text, *, size=14, bold=False, color=COLOR_DARK,
             align=PP_ALIGN.LEFT, font="Calibri", italic=False):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = color; r.font.name = font
    return tb

def add_bullets(s, x, y, w, h, items, *, size=14, color=COLOR_DARK, bullet_color=COLOR_ACCENT, gap=6):
    tb = s.shapes.add_textbox(x, y, w, h); tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        sub = item.startswith("  ")
        text = item.strip()
        r1 = p.add_run(); r1.text = ("–  " if sub else "•  ")
        r1.font.size = Pt(size); r1.font.bold = True; r1.font.color.rgb = bullet_color; r1.font.name = "Calibri"
        if sub: p.level = 1
        # suporte a **negrito** inline simples
        for j, chunk in enumerate(text.split("**")):
            if chunk == "": continue
            r = p.add_run(); r.text = chunk
            r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = "Calibri"
            r.font.bold = (j % 2 == 1)
    return tb

def header(s, title, subtitle=None):
    add_rect(s, Emu(0), Emu(0), Inches(0.25), SH, COLOR_PRIMARY)
    add_text(s, Inches(0.55), Inches(0.30), Inches(12.4), Inches(0.6), title,
             size=26, bold=True, color=COLOR_PRIMARY)
    if subtitle:
        add_text(s, Inches(0.55), Inches(0.86), Inches(12.4), Inches(0.4), subtitle,
                 size=13, color=COLOR_MUTED)
    add_rect(s, Inches(0.55), Inches(1.26), Inches(1.4), Inches(0.05), COLOR_ACCENT)

def footer(s, n):
    add_text(s, Inches(0.55), Inches(7.12), Inches(9.5), Inches(0.3),
             "TP4 DemoDay · Grupo 02 · Efeito Dominó no Mercado de Transferências · UFMG",
             size=9, color=COLOR_MUTED)
    add_text(s, Inches(11.6), Inches(7.12), Inches(1.4), Inches(0.3), f"{n} / 14",
             size=9, color=COLOR_MUTED, align=PP_ALIGN.RIGHT)

def add_table(s, x, y, w, h, data, *, highlight_rows=(), col_widths=None, font_size=11):
    rows, cols = len(data), len(data[0])
    t = s.shapes.add_table(rows, cols, x, y, w, h).table
    if col_widths:
        for i, cw in enumerate(col_widths): t.columns[i].width = cw
    for r in range(rows):
        for c in range(cols):
            cell = t.cell(r, c); cell.text = ""
            cell.margin_left = cell.margin_right = Inches(0.07)
            cell.margin_top = cell.margin_bottom = Inches(0.03)
            p = cell.text_frame.paragraphs[0]; run = p.add_run(); run.text = str(data[r][c])
            run.font.name = "Calibri"
            if r == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = COLOR_PRIMARY
                run.font.color.rgb = COLOR_WHITE; run.font.bold = True; run.font.size = Pt(font_size)
            else:
                cell.fill.solid()
                if r in highlight_rows:
                    cell.fill.fore_color.rgb = RGBColor(0xE8, 0xF5, 0xE9); run.font.bold = True
                else:
                    cell.fill.fore_color.rgb = COLOR_WHITE if r % 2 else COLOR_LIGHT
                run.font.color.rgb = COLOR_DARK; run.font.size = Pt(font_size)
            if c > 0: p.alignment = PP_ALIGN.CENTER
    return t

def kpi_card(s, x, y, w, h, kpi, label, color):
    add_rect(s, x, y, w, h, COLOR_LIGHT)
    add_rect(s, x, y, w, Inches(0.10), color)
    add_text(s, x, y + Inches(0.42), w, Inches(1.0), kpi, size=40, bold=True,
             color=color, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.15), y + Inches(1.55), w - Inches(0.3), Inches(0.8), label,
             size=11.5, color=COLOR_DARK, align=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 1 — Capa
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, Emu(0), Emu(0), SW, SH, COLOR_PRIMARY)
add_rect(s, Inches(0), Inches(6.7), SW, Inches(0.8), COLOR_DARK)
add_text(s, Inches(0.8), Inches(1.5), Inches(11.5), Inches(0.5),
         "TP4 · DemoDay", size=18, color=COLOR_ACCENT, bold=True)
add_text(s, Inches(0.8), Inches(2.1), Inches(11.7), Inches(2.0),
         "Efeito Dominó no Mercado\nde Transferências de Futebol", size=44, bold=True, color=COLOR_WHITE)
add_text(s, Inches(0.8), Inches(4.35), Inches(11.5), Inches(0.5),
         "Existe um \"prêmio do vendedor\"? Uma abordagem causal (Double ML)",
         size=19, color=COLOR_LIGHT)
add_text(s, Inches(0.8), Inches(5.45), Inches(11.5), Inches(0.4),
         "Grupo 02  ·  Ciência de Dados Aplicada ao Futebol  ·  UFMG", size=13, color=COLOR_LIGHT)
add_text(s, Inches(0.8), Inches(5.82), Inches(11.5), Inches(0.4),
         "Carlos · César · Letícia · Lucas · Lucca", size=13, color=COLOR_LIGHT)
add_text(s, Inches(0.8), Inches(6.9), Inches(11.5), Inches(0.4), "25 / 06 / 2026",
         size=11, color=COLOR_MUTED)

# ============================================================
# SLIDE 2 — Contexto & Problema
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "Contexto & Problema", "O mercado de transferências como rede de dependências")
add_bullets(s, Inches(0.7), Inches(1.6), Inches(12.0), Inches(2.2), [
    "Quando um clube faz uma **grande venda**, dois sinais chegam ao mercado ao mesmo tempo: "
    "ele tem **caixa** e precisa **repor** o jogador. Os vendedores sabem disso.",
    "Intuição — o **\"efeito dominó\"**: esse clube perde poder de barganha e passa a pagar "
    "**acima do valor justo** nas compras seguintes.",
], size=15, gap=10)
add_rect(s, Inches(0.7), Inches(4.05), Inches(12.0), Inches(1.5), COLOR_LIGHT)
add_rect(s, Inches(0.7), Inches(4.05), Inches(0.12), Inches(1.5), COLOR_RED)
add_text(s, Inches(1.0), Inches(4.2), Inches(11.5), Inches(0.4), "A armadilha óbvia",
         size=13, bold=True, color=COLOR_RED)
add_text(s, Inches(1.0), Inches(4.6), Inches(11.5), Inches(0.9),
         "Clubes ricos vendem caro E compram caro ao mesmo tempo. Como saber se o que vemos é o "
         "prêmio do vendedor — ou apenas o porte financeiro do clube? Separar as duas coisas é o "
         "desafio central do trabalho.", size=14, color=COLOR_DARK)
footer(s, 2)

# ============================================================
# SLIDE 3 — Pergunta & Hipótese
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "Pergunta & Hipótese", "Causal, não correlacional")
add_rect(s, Inches(0.7), Inches(1.7), Inches(12.0), Inches(1.7), COLOR_LIGHT)
add_rect(s, Inches(0.7), Inches(1.7), Inches(0.12), Inches(1.7), COLOR_ACCENT)
add_text(s, Inches(1.0), Inches(1.85), Inches(11.4), Inches(0.4), "Pergunta operacional",
         size=12, bold=True, color=COLOR_ACCENT)
add_text(s, Inches(1.0), Inches(2.25), Inches(11.4), Inches(1.1),
         "A liquidez recente do comprador (receita de vendas na temporada) tem efeito CAUSAL sobre "
         "o sobrepreço que ele paga nas contratações seguintes, depois de neutralizar os "
         "confundidores estruturais?", size=16, color=COLOR_DARK)
add_text(s, Inches(0.7), Inches(3.75), Inches(12.0), Inches(0.4), "Hipótese inicial (conjectura)",
         size=14, bold=True, color=COLOR_PRIMARY)
add_bullets(s, Inches(0.7), Inches(4.15), Inches(12.0), Inches(1.6), [
    "Sim — e o mecanismo seria de **sinalização**, não de volume de dinheiro: o clube pagaria mais "
    "porque o mercado sabe que ele **precisa comprar**.",
    "  Guardem essa conjectura: no Resultado 3 mostramos que os dados **não a sustentam de forma "
    "robusta** — exemplo de rigor, não de fraqueza.",
], size=14, gap=8)
footer(s, 3)

# ============================================================
# SLIDE 4 — A Armadilha Causal (confundidores)
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "A Armadilha Causal", "6 confundidores mapeados antes de modelar")
conf = [
    ["#", "Confundidor", "Como tratamos"],
    ["C1", "Clube rico (vende caro E compra caro)", "Proxies de rede (PageRank) + elenco em W"],
    ["C2", "Causalidade reversa (compra → venda)", "Teste de tratamento defasado (t−1)"],
    ["C3", "Causa comum (UCL, troca de técnico)", "Não controlado (sem dados esportivos)"],
    ["C4", "Sazonalidade (deadline day)", "Dummies de temporada (Etapa 1 e W)"],
    ["C5", "Inflação de mercado entre anos", "Dummies de temporada absorvem inflação"],
    ["C6", "Viés de seleção do grupo tratado", "Proxies de rede + sensibilidade ao corte"],
]
add_table(s, Inches(0.7), Inches(1.55), Inches(12.0), Inches(3.6), conf,
          highlight_rows=(1,), col_widths=[Inches(0.7), Inches(5.6), Inches(5.7)], font_size=12)
add_rect(s, Inches(0.7), Inches(5.35), Inches(12.0), Inches(1.1), COLOR_LIGHT)
add_text(s, Inches(0.9), Inches(5.46), Inches(11.6), Inches(0.3),
         "Por que isto é o coração do projeto", size=12, bold=True, color=COLOR_PRIMARY)
add_text(s, Inches(0.9), Inches(5.78), Inches(11.6), Inches(0.6),
         "Sem tratar C1–C6, qualquer correlação observada seria espúria. O destaque é o C1: é ele "
         "que faz a correlação ingênua parecer um efeito causal — como veremos no Resultado 1.",
         size=12, color=COLOR_DARK)
footer(s, 4)

# ============================================================
# SLIDE 5 — Metodologia: 2 etapas
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "Metodologia", "Duas etapas encadeadas (estado da arte — Chernozhukov et al., 2018)")

def stage_card(x, y, w, h, num, title, lines, color):
    add_rect(s, x, y, w, h, COLOR_LIGHT)
    add_rect(s, x, y, w, Inches(0.5), color)
    add_text(s, x + Inches(0.25), y + Inches(0.08), w, Inches(0.4),
             f"ETAPA {num}  ·  {title}", size=13, bold=True, color=COLOR_WHITE)
    tb = s.shapes.add_textbox(x + Inches(0.25), y + Inches(0.66), w - Inches(0.5), h - Inches(0.7))
    tf = tb.text_frame; tf.word_wrap = True
    for i, (lab, val) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph(); p.space_after = Pt(4)
        r1 = p.add_run(); r1.text = lab; r1.font.size = Pt(12); r1.font.bold = True
        r1.font.color.rgb = COLOR_PRIMARY; r1.font.name = "Calibri"
        if val:
            r2 = p.add_run(); r2.text = " " + val; r2.font.size = Pt(12)
            r2.font.color.rgb = COLOR_DARK; r2.font.name = "Calibri"

stage_card(Inches(0.7), Inches(1.6), Inches(5.9), Inches(4.3), 1, "Modelo Hedônico", [
    ("Objetivo:", "estimar o \"preço justo\" do jogador (ML)."),
    ("Target:", "ln(fee) — log do valor da transação."),
    ("Features:", "idade, idade², log(market value),"),
    ("", "posição, liga, temporada."),
    ("Saída:", "resíduo Y = ln(P) − ln(P̂)"),
    ("", "→ o \"sobrepreço\" puro."),
], COLOR_PRIMARY)
stage_card(Inches(6.85), Inches(1.6), Inches(5.9), Inches(4.3), 2, "Double ML (causal)", [
    ("Objetivo:", "isolar o efeito causal da liquidez"),
    ("", "sobre o sobrepreço."),
    ("D (tratamento):", "log da receita de vendas."),
    ("Y (resultado):", "resíduo hedônico (Etapa 1)."),
    ("W (confund.):", "rede + elenco + liga + temporada (19)."),
    ("Técnica:", "ortogonalização de Neyman + cross-fitting."),
], COLOR_ACCENT)
arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.55), Inches(3.5), Inches(0.35), Inches(0.5))
arrow.fill.solid(); arrow.fill.fore_color.rgb = COLOR_DARK; arrow.line.fill.background()
add_text(s, Inches(0.7), Inches(6.15), Inches(12.0), Inches(0.9),
         "A separação garante que o resíduo Y NÃO contenha o comportamento financeiro do clube — "
         "esse fica reservado para os confundidores W da Etapa 2.", size=12, color=COLOR_DARK)
footer(s, 5)

# ============================================================
# SLIDE 6 — Etapa 1: Hedônico
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "Etapa 1 · Modelo Hedônico", "Funil de dados + 4 modelos + interpretabilidade")
# Funil
for i, (txt, col) in enumerate([("44.627 movimentações brutas", COLOR_MUTED),
                                 ("filtro: pagas, ≥ €250k, MV válido", COLOR_ACCENT),
                                 ("5.146 transferências · 8 temp · 7 ligas", COLOR_PRIMARY)]):
    y = Inches(1.65 + i*0.62)
    add_rect(s, Inches(0.7), y, Inches(5.6 - i*0.5), Inches(0.5), col)
    add_text(s, Inches(0.85), y + Inches(0.09), Inches(5.2), Inches(0.4), txt,
             size=11.5, bold=True, color=COLOR_WHITE)
add_text(s, Inches(0.7), Inches(3.7), Inches(5.7), Inches(0.4), "Ranking SHAP (impacto)",
         size=12, bold=True, color=COLOR_PRIMARY)
for i, (name, val) in enumerate([("log_mv", 1.0), ("age / age²", 0.5), ("liga", 0.3),
                                 ("posição", 0.18), ("temporada", 0.12)]):
    y = Inches(4.12 + i*0.46)
    add_text(s, Inches(0.7), y, Inches(1.5), Inches(0.4), name, size=11, bold=True, color=COLOR_DARK)
    add_rect(s, Inches(2.2), y + Inches(0.06), Inches(3.0*val), Inches(0.26), COLOR_ACCENT)
# Tabela de modelos
tbl = [
    ["Modelo", "CV R²", "Test RMSE", "Test R²"],
    ["XGBoost", "0,722", "0,694", "0,746"],
    ["LightGBM", "0,729", "0,683", "0,755"],
    ["Random Forest", "0,731", "0,673", "0,762"],
    ["SVR (RBF)", "0,708", "0,716", "0,730"],
]
add_table(s, Inches(6.8), Inches(1.65), Inches(5.9), Inches(2.4), tbl,
          highlight_rows=(3,), col_widths=[Inches(2.3), Inches(1.2), Inches(1.2), Inches(1.2)], font_size=11)
add_rect(s, Inches(6.8), Inches(4.2), Inches(5.9), Inches(0.55), RGBColor(0xE8, 0xF5, 0xE9))
add_text(s, Inches(6.95), Inches(4.3), Inches(5.7), Inches(0.4),
         "✓ Random Forest — R² 0,76 no holdout de 2025", size=12.5, bold=True, color=COLOR_GREEN)
add_bullets(s, Inches(6.8), Inches(4.95), Inches(5.9), Inches(1.8), [
    "Split **temporal**: treina 2017–2024, testa 2025 (sem vazamento).",
    "Resíduo gerado **out-of-fold** (cross_val_predict) — sem vazamento in-sample.",
    "log_mv domina; nenhuma feature do clube comprador entra (por design).",
], size=11.5, gap=5)
footer(s, 6)

# ============================================================
# SLIDE 7 — Etapa 2: Double ML
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "Etapa 2 · Double Machine Learning", "Como isolamos o efeito causal")
for i, (tag, txt, col) in enumerate([
    ("D", "Tratamento: log da receita de vendas do clube na temporada (liquidez recente).", COLOR_PRIMARY),
    ("Y", "Resultado: o resíduo hedônico da Etapa 1 — o sobrepreço acima do valor justo.", COLOR_ACCENT),
    ("W", "Confundidores (19): PageRank e graus na rede, tamanho de elenco, 7 dummies de "
          "temporada, 6 de liga. Controla C1, C4, C5.", COLOR_GREEN)]):
    y = Inches(1.65 + i*1.0)
    add_rect(s, Inches(0.7), y, Inches(0.7), Inches(0.7), col)
    add_text(s, Inches(0.7), y + Inches(0.13), Inches(0.7), Inches(0.5), tag,
             size=22, bold=True, color=COLOR_WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1.6), y + Inches(0.04), Inches(11.0), Inches(0.95), txt, size=14, color=COLOR_DARK)
add_rect(s, Inches(0.7), Inches(4.8), Inches(12.0), Inches(1.55), COLOR_LIGHT)
add_text(s, Inches(0.9), Inches(4.92), Inches(11.6), Inches(0.4),
         "Ortogonalização de Neyman + cross-fitting (5 folds)", size=13, bold=True, color=COLOR_PRIMARY)
add_text(s, Inches(0.9), Inches(5.32), Inches(11.6), Inches(1.0),
         "Tiramos de Y e de D tudo o que os confundidores W explicam (Ỹ, D̃) e estimamos θ na "
         "variação que sobra — exógena por construção. Erros-padrão HC1 e CLUSTERIZADOS por "
         "clube×temporada (974 clusters), já que D varia nesse nível, não por transferência.",
         size=13, color=COLOR_DARK)
footer(s, 7)

# ============================================================
# SLIDE 8 — Resultado 1: A Reviravolta
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "Resultado 1 · A Reviravolta", "A correlação mente")
s.shapes.add_picture(IMG_PROG, Inches(0.6), Inches(1.55), height=Inches(4.2))
add_text(s, Inches(7.6), Inches(1.7), Inches(5.1), Inches(0.5), "O que aconteceu", size=14, bold=True, color=COLOR_PRIMARY)
add_bullets(s, Inches(7.6), Inches(2.2), Inches(5.2), Inches(3.5), [
    "**OLS ingênuo:** θ = 0,0142* — significativo. \"Prova\" o prêmio do vendedor.",
    "**Double ML:** θ = 0,0044 — IC cruza zero (HC1 **e** cluster). **Não significativo.**",
    "**Placebo** (D embaralhado): θ ≈ 0,002 — quase do tamanho do efeito real.",
    "A correlação bruta era, em grande parte, o **confundidor 'clube rico'**.",
], size=13, gap=9)
add_rect(s, Inches(7.6), Inches(5.7), Inches(5.2), Inches(0.9), RGBColor(0xFB, 0xEE, 0xE6))
add_text(s, Inches(7.75), Inches(5.8), Inches(5.0), Inches(0.7),
         "Sem o DML, teríamos publicado um efeito que, na média de 8 temporadas, não existe.",
         size=12, bold=True, color=COLOR_ACCENT)
footer(s, 8)

# ============================================================
# SLIDE 9 — Resultado 2: Onde o efeito vive
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "Resultado 2 · Onde o Efeito Vive", "O achado central: condicional ao regime de mercado")
s.shapes.add_picture(IMG_SAZ, Inches(0.6), Inches(1.55), height=Inches(4.3))
add_bullets(s, Inches(8.7), Inches(1.7), Inches(4.1), Inches(4.0), [
    "Na **média**, efeito nulo. Mas por temporada a história muda.",
    "**2022 (+6,4%) e 2023 (+4,7%)** — únicos significativos, e **sobrevivem a Bonferroni/FDR**.",
    "Anos do **boom pós-COVID**: escassez de talento + urgência de reposição.",
    "Demais 6 anos: indistinguível de zero.",
], size=13, gap=9)
add_rect(s, Inches(8.7), Inches(5.75), Inches(4.1), Inches(0.85), RGBColor(0xE8, 0xF5, 0xE9))
add_text(s, Inches(8.85), Inches(5.85), Inches(3.9), Inches(0.7),
         "O prêmio do vendedor não é lei universal — emerge em mercado aquecido.",
         size=12, bold=True, color=COLOR_GREEN)
footer(s, 9)

# ============================================================
# SLIDE 10 — Resultado 3: O Mecanismo
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "Resultado 3 · O Mecanismo", "Honestidade: a conjectura inicial não se sustenta")
tbl = [
    ["Tratamento", "θ (bruto)", "Com controle de 'clube ativo'"],
    ["D₁ = log(receita €)", "0,004  n.s.", "—"],
    ["D₂ = log(nº de vendas)", "0,049 *", "0,015  (p=0,40)  n.s."],
    ["D₃ = venda blockbuster >€30M", "0,075 *", "0,005  (p=0,86)  n.s."],
]
add_table(s, Inches(0.7), Inches(1.6), Inches(12.0), Inches(2.0), tbl,
          col_widths=[Inches(4.4), Inches(2.8), Inches(4.8)], font_size=12.5)
add_bullets(s, Inches(0.7), Inches(3.9), Inches(12.0), Inches(2.0), [
    "No bruto, D₂/D₃ eram significativos → **sugeria** que o que importa é sinalizar necessidade.",
    "**MAS** ao controlar por \"clube de janela movimentada\" (nº de compras + gasto), **ambos "
    "perdem toda a significância**. O sinal era em boa parte \"clube ativo compra e vende muito\".",
], size=14, gap=10)
add_rect(s, Inches(0.7), Inches(5.7), Inches(12.0), Inches(0.95), RGBColor(0xFB, 0xEE, 0xE6))
add_text(s, Inches(0.9), Inches(5.82), Inches(11.6), Inches(0.75),
         "Conclusão: o \"mecanismo de sinalização\" é hipótese EXPLORATÓRIA, não achado confirmado. "
         "Os resultados sólidos são o Achado 1 (a correlação mente) e o efeito sazonal 2022–2023.",
         size=12.5, bold=True, color=COLOR_ACCENT)
footer(s, 10)

# ============================================================
# SLIDE 11 — Resultado 4: Robustez
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "Resultado 4 · Robustez", "Blindando o achado de 2022–2023")
cards = [
    ("Placebo dentro\nda safra", "p = 0,005", "permuta D em 2022/2023 → θ colapsa para ~0"),
    ("Sensibilidade\nao corte de fee", "€0–€1M", "ATE nulo em todos; 2022/2023 estáveis"),
    ("Robustness\nValue", "24–27%", "força que um confundidor omitido precisaria ter"),
    ("Tratamento\ndefasado (C2)", "não concl.", "lag enfraquece causalidade reversa"),
]
for i, (title, big, desc) in enumerate(cards):
    x = Inches(0.7 + i*3.06)
    add_rect(s, x, Inches(1.65), Inches(2.86), Inches(2.5), COLOR_LIGHT)
    add_rect(s, x, Inches(1.65), Inches(2.86), Inches(0.10), COLOR_PRIMARY)
    add_text(s, x + Inches(0.12), Inches(1.82), Inches(2.62), Inches(0.7), title,
             size=12.5, bold=True, color=COLOR_PRIMARY, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.12), Inches(2.55), Inches(2.62), Inches(0.5), big,
             size=20, bold=True, color=COLOR_ACCENT, align=PP_ALIGN.CENTER)
    add_text(s, x + Inches(0.12), Inches(3.25), Inches(2.62), Inches(0.85), desc,
             size=10.5, color=COLOR_DARK, align=PP_ALIGN.CENTER)
add_rect(s, Inches(0.7), Inches(4.5), Inches(12.0), Inches(1.9), COLOR_LIGHT)
add_text(s, Inches(0.9), Inches(4.62), Inches(11.6), Inches(0.4),
         "Inferência levada a sério", size=13, bold=True, color=COLOR_PRIMARY)
add_bullets(s, Inches(0.9), Inches(5.05), Inches(11.6), Inches(1.3), [
    "Erros-padrão **clusterizados** por clube×temporada (974 clusters) — não tratamos 5.146 obs como independentes.",
    "**Correção de múltiplas comparações** (Bonferroni + FDR) nos 8 testes sazonais — 2022/2023 sobrevivem.",
    "Placebos globais (D embaralhado, ruído) ≈ 0 — o modelo não inventa efeito.",
], size=12, gap=5)
footer(s, 11)

# ============================================================
# SLIDE 12 — Conclusões & Implicações
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "Conclusões & Implicações", "O prêmio existe — sob condições")
kpi_card(s, Inches(0.7), Inches(1.6), Inches(3.9), Inches(2.2), "Nulo",
         "efeito MÉDIO em 8 temporadas\n(a correlação ingênua mentia)", COLOR_PRIMARY)
kpi_card(s, Inches(4.75), Inches(1.6), Inches(3.9), Inches(2.2), "2022–23",
         "efeito significativo só no\nboom pós-COVID (robusto)", COLOR_GREEN)
kpi_card(s, Inches(8.8), Inches(1.6), Inches(3.9), Inches(2.2), "Regime",
         "o timing importa mais que\no volume de caixa", COLOR_ACCENT)
add_text(s, Inches(0.7), Inches(4.05), Inches(12.0), Inches(0.4), "Implicações práticas",
         size=14, bold=True, color=COLOR_PRIMARY)
add_bullets(s, Inches(0.7), Inches(4.45), Inches(12.0), Inches(2.2), [
    "**Timing de compra:** em mercado aquecido, comprar logo após uma grande venda é caro. Quem pode esperar leva vantagem.",
    "**Inteligência de mercado:** o IVB aponta clubes mais vulneráveis ao prêmio (ilustrativo — base para scouting financeiro).",
    "**Alerta metodológico:** o OLS ingênuo levaria à conclusão errada. Correlação não é causalidade.",
], size=13.5, gap=9)
footer(s, 12)

# ============================================================
# SLIDE 13 — Limitações & Trabalho Futuro
# ============================================================
s = prs.slides.add_slide(BLANK)
header(s, "Limitações & Trabalho Futuro", "O que reconhecemos honestamente")
add_bullets(s, Inches(0.7), Inches(1.7), Inches(12.0), Inches(3.0), [
    "**Granularidade sazonal:** o tratamento é por temporada, não por dia. O efeito da janela de 30 dias está diluído → estimativa conservadora.",
    "**Identificação por observáveis:** C1 controlado por proxies (sem bloco financeiro direto, que seria bad control). O Robustness Value mitiga, mas não substitui um instrumento.",
    "**Sem performance individual:** não integramos gols/xG/assistências (FBref).",
    "**IVB via R-learner:** ranking ilustrativo; CausalForestDML (econml) daria resultado mais robusto.",
], size=14, gap=11)
add_rect(s, Inches(0.7), Inches(5.5), Inches(12.0), Inches(1.0), COLOR_LIGHT)
add_text(s, Inches(0.9), Inches(5.62), Inches(11.6), Inches(0.4), "Trabalho futuro",
         size=12.5, bold=True, color=COLOR_PRIMARY)
add_text(s, Inches(0.9), Inches(5.98), Inches(11.6), Inches(0.5),
         "Datas diárias (lag de curto prazo conclusivo) · PSM · curva de decaimento θ(Δt) · "
         "análise do lado vendedor (\"clubes predadores\") · instrumento para C1.",
         size=12.5, color=COLOR_DARK)
footer(s, 13)

# ============================================================
# SLIDE 14 — Encerramento
# ============================================================
s = prs.slides.add_slide(BLANK)
add_rect(s, Emu(0), Emu(0), SW, SH, COLOR_PRIMARY)
add_rect(s, Inches(0), Inches(6.7), SW, Inches(0.8), COLOR_DARK)
add_text(s, Inches(1.2), Inches(2.3), Inches(11.0), Inches(2.0),
         "\"O prêmio do vendedor não é uma lei do mercado — ele emerge nas janelas de mercado "
         "aquecido, e some no restante.\"", size=26, bold=True, color=COLOR_WHITE, italic=True)
add_text(s, Inches(1.2), Inches(4.7), Inches(11.0), Inches(0.5),
         "Obrigado!  ·  Estamos à disposição para perguntas.", size=18, color=COLOR_ACCENT, bold=True)
add_text(s, Inches(1.2), Inches(5.5), Inches(11.0), Inches(0.4),
         "Grupo 02 · Ciência de Dados Aplicada ao Futebol · UFMG", size=13, color=COLOR_LIGHT)

# ---------- Salva ----------
out = "/Users/lucaspedras/TP-CDAF/output/TP4_Apresentacao_Final_Grupo02.pptx"
prs.save(out)
print(f"OK: {out} ({len(prs.slides._sldIdLst)} slides)")
